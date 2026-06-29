#!/usr/bin/env python3
"""Generate Hindi-only EWUMS UJS department presentation (.pptx) with PRP logo."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# PRP brand colors
NAVY = RGBColor(0x0F, 0x17, 0x2A)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
FOOTER_TEXT = "PRP Geospatial Solutions | ewumsujs.com"
FONT = "Nirmala UI"

SCRIPT_DIR = Path(__file__).parent
REPO_ROOT = SCRIPT_DIR.parent.parent
LOGO_CANDIDATES = [
    REPO_ROOT / "frontend" / "web" / "public" / "prp-logo.png",
    REPO_ROOT / "public" / "prp-logo.png",
]
OUTPUT = SCRIPT_DIR / "EWUMS-UJS-Presentation-Hindi.pptx"

SLIDES = [
    {
        "type": "title",
        "title": "उत्तराखंड जल संस्थान के लिए\nएकीकृत जल उपयोगिता प्रबंधन",
        "subtitle": "एकीकृत जल उपयोगिता प्रबंधन प्रणाली — UJS",
        "tagline": "EWUMS — उद्यम जल उपयोगिता प्रबंधन प्रणाली (S2T2R)",
        "notes": (
            "सुप्रभात। आज हम EWUMS प्रस्तुत कर रहे हैं — UJS के लिए योजना, निर्माण, GIS, "
            "संचालन, बिलिंग और कार्यकारी देखरेख को कवर करने वाला एक एकीकृत डिजिटल मंच।"
        ),
    },
    {
        "type": "bullets",
        "title": "UJS को एकीकृत मंच की आवश्यकता क्यों है",
        "hindi": "चुनौती — बिखरी हुई जानकारी और मैन्युअल प्रक्रियाएँ",
        "bullets": [
            "योजनाएँ स्प्रेडशीट, फ़ाइलों और अलग-अलग GIS उपकरणों में ट्रैक होती हैं",
            "भौतिक बनाम वित्तीय प्रगति पर देरी से दृश्यता",
            "मंडल-वार डेटा सिलो (गढ़वाल / कुमाऊँ सर्कल)",
            "उपभोक्ता बिलिंग, शिकायतें और O&M मानचित्र पर संपत्तियों से जुड़े नहीं",
            "ऑडिट और अनुमोदन ट्रेल पुनर्निर्माण करना कठिन",
        ],
        "notes": (
            "अधिकारी योजनाओं का प्रबंधन करने के बजाय डेटा का मिलान करने में समय व्यतीत करते हैं। "
            "EWUMS इसे अंत से अंत तक संबोधित करता है।"
        ),
    },
    {
        "type": "bullets",
        "title": "EWUMS — एक मंच, पूर्ण जीवनचक्र",
        "hindi": "समाधान — योजना से O&M तक एक ही प्रणाली",
        "bullets": [
            "21 एकीकृत मॉड्यूल — 17 आज ewumsujs.com पर लाइव",
            "DPR → निर्माण → MB → RA बिल → अंतिम बिल → GIS हस्तांतरण → O&M → बिलिंग",
            "UJS क्षेत्र संरचना के अनुरूप मंडल-स्कोप पहुँच (Migration 084 मंडल)",
            "जिला सीमाओं के साथ वेब GIS (देहरादून, हरिद्वार, राज्यव्यापी)",
            "मोबाइल बिलिंग और उपभोक्ता पोर्टल (जल मित्र)",
        ],
        "notes": "",
    },
    {
        "type": "table",
        "title": "UJS संगठन संरचना के लिए निर्मित",
        "headers": ["स्तर", "भूमिकाएँ", "मंच पहुँच"],
        "rows": [
            ["राज्य मुख्यालय (देहरादून)", "MD, CGM, CE", "राज्यव्यापी डैशबोर्ड और सभी मंडल"],
            ["सर्कल", "SE गढ़वाल / कुमाऊँ", "सर्कल-व्यापी योजनाएँ"],
            ["मंडल", "EE, AE, JE", "मंडल-स्कोप परियोजनाएँ और GIS"],
            ["क्षेत्र", "ठेकेदार, लेखा", "निर्माण वर्कफ़्लो और भुगतान"],
            ["नागरिक", "उपभोक्ता", "जल मित्र पोर्टल — बिल, शिकायतें"],
        ],
        "footer_note": (
            "संदर्भ मंडल: देहरादून उत्तर (DIV-DNN), हरिद्वार ग्रामीण (DIV-HRR), "
            "करणप्रयाग (DIV-KPG), नैनीताल (DIV-NTL)"
        ),
        "notes": "",
    },
    {
        "type": "numbered",
        "title": "योजना एवं निर्माण — लाइव मॉड्यूल",
        "hindi": "योजना एवं निर्माण",
        "items": [
            "DPR और योजना प्रबंधन — 12-चरण अनुमोदन",
            "निर्माण प्रबंधन — प्रगति डैशबोर्ड",
            "MB प्रबंधन — माप और सत्यापन",
            "BOQ और मात्रा समाधान",
            "RA बिल और ठेकेदार भुगतान",
            "अंतिम बिल और परियोजना समापन",
            "भूमि अधिग्रहण प्रबंधन — GIS पार्सल ट्रेस",
        ],
        "notes": (
            "प्रस्ताव से अंतिम बिल तक, प्रत्येक चरण भूमिका-आधारित अनुमोदन के साथ दर्ज "
            "(JE → AE → EE → वित्त)।"
        ),
    },
    {
        "type": "sections",
        "title": "GIS, संचालन और राजस्व — लाइव मॉड्यूल",
        "sections": [
            ("GIS और संपत्ति", "वेब मानचित्र एक्सप्लोरर, संपत्ति रजिस्ट्री, फ़ीचर क्लास"),
            ("संचालन", "O&M वर्कफ़्लो (14 चरण), SCADA/IoT, जल गुणवत्ता, शिकायतें, उपभोक्ता"),
            ("वाणिज्यिक", "बिलिंग और राजस्व, ERP/GL एकीकरण, 15.12 रिपोर्ट"),
            ("बुद्धिमत्ता", "कार्यकारी डैशबोर्ड, संपत्ति जीवनचक्र, एनालिटिक्स"),
        ],
        "footer_note": (
            "आंशिक (रोडमैप): मोबाइल कार्यबल विस्तार, एंटरप्राइज DMS, "
            "QR स्कैन ऐप, ML predictive maintenance"
        ),
        "notes": "",
    },
    {
        "type": "bullets",
        "title": "स्थान — सामान्य धागा",
        "hindi": "GIS — हर संपत्ति और योजना मानचित्र पर",
        "bullets": [
            "जिला-लॉक मानचित्र अधिकार क्षेत्र (देहरादून, हरिद्वार)",
            "पाइपलाइन, जलाशय, FHTC परतें परियोजनाओं से जुड़ी",
            "भूमि अधिग्रहण: स्वचालित संरेखण ट्रेस + कैडास्ट्रल इंटरसेक्ट",
            "राजस्व GIS एनालिटिक्स — भूगोल के अनुसार संग्रह",
            "क्षेत्र सत्यापन के लिए उपग्रह बेसमैप (Google / Esri)",
        ],
        "footer_note": "डेमो: Map Explorer → देहरादून जिला → योजना संपत्तियाँ",
        "notes": "",
    },
    {
        "type": "demo",
        "title": "डेमो 1 — कमांड सेंटर",
        "hindi": "कार्यकारी डैशबोर्ड",
        "login": "admin@egip.local / Admin@123",
        "navigate": "साइडबार → Executive Dashboard (/dashboard)",
        "show": [
            "पाँच KPI कार्ड — संपत्ति, अलर्ट, महत्वपूर्ण संपत्ति, परियोजना प्रगति, स्वास्थ्य",
            "परियोजना प्रगति चार्ट — क्लेमेंट टाउन और बहादराबाद योजनाएँ",
            "महत्वपूर्ण संपत्ति तालिका",
            "20 मंच मॉड्यूल का लिंक",
        ],
        "script": (
            "एक स्क्रीन से, नेतृत्व राज्यव्यापी स्वास्थ्य देखता है। "
            "प्रत्येक KPI PostgreSQL से लाइव है — स्थिर रिपोर्ट नहीं।"
        ),
        "notes": "लाइव डेमो — 2 मिनट। demo-seed.sql लागू करके अभ्यास करें।",
    },
    {
        "type": "demo",
        "title": "डेमो 2 — योजना कार्यान्वयन",
        "hindi": "निर्माण प्रवाह",
        "login": "je.dnn@egip.local / JE@123 (देहरादून उत्तर JE)",
        "navigate": "Projects → Clement Town Water Supply Scheme",
        "show": [
            "निर्माण डैशबोर्ड — भौतिक 72% / वित्तीय 65%",
            "MB प्रविष्टि (माप पुस्तक)",
            "BOQ समाधान",
            "RA बिल AE अनुमोदन लंबित",
            "वर्कफ़्लो इनबॉक्स — AE के रूप में अनुमोदित करें (ae.dnn@egip.local)",
        ],
        "script": (
            "मंडल कर्मचारी केवल अपनी योजनाएँ देखते हैं। "
            "JE मापता है, AE जाँचता है, EE अनुमोदित करता है — पूर्ण ऑडिट ट्रेल।"
        ),
        "notes": "लाइव डेमो — 4 मिनट।",
    },
    {
        "type": "demo",
        "title": "डेमो 3 — राजस्व और नागरिक सेवा",
        "hindi": "बिलिंग एवं शिकायत",
        "login": "accounts.hrr@egip.local / Accounts@123 (हरिद्वार ग्रामीण)",
        "navigate": "Billing → Revenue KPIs",
        "show": [
            "उपभोक्ता कनेक्शन और संग्रह दक्षता",
            "बहादराबाद उपभोक्ता के लिए नमूना बिल",
            "शिकायत रजिस्टर (यदि सीड किया गया)",
            "वैकल्पिक: जल मित्र उपभोक्ता पोर्टल — FHTC-HRR-DEMO-001",
        ],
        "script": (
            "कमीशनिंग के बाद, वही मंच टैरिफ, मीटर रीडिंग और "
            "उपभोक्ता स्व-सेवा का प्रबंधन करता है।"
        ),
        "notes": "लाइव डेमो — 4 मिनट।",
    },
    {
        "type": "bullets",
        "title": "क्षेत्र संचालन — मोबाइल तैयार",
        "hindi": "मोबाइल — क्षेत्र बिलिंग और GPS कैप्चर",
        "bullets": [
            "मोबाइल बिलिंग मार्ग (/mobile-billing) — मीटर रीडिंग, ऑफ़लाइन कतार",
            "रीडिंग पर GPS और फ़ोटो साक्ष्य",
            "उपभोक्ता पोर्टल — OTP लॉगिन, हिंदी/अंग्रेज़ी/गढ़वाली/कुमाऊँ (जल मित्र)",
            "QR संपत्ति पहचान (आंशिक — स्कैन API लाइव, समर्पित ऐप रोडमैप)",
        ],
        "notes": "",
    },
    {
        "type": "bullets",
        "title": "विश्वास और अनुपालन",
        "hindi": "सुरक्षा, ऑडिट ट्रेल और मंडल पहुँच नियंत्रण",
        "bullets": [
            "भूमिका-आधारित पहुँच — 25+ UJS भूमिकाएँ (MD, SE, EE, JE, बिलिंग अधिकारी, आदि)",
            "मंडल अलगाव — हरिद्वार JE देहरादून उत्तर डेटा नहीं देख सकता",
            "पूर्ण ऑडिट लॉग — उपयोगकर्ता क्रियाएँ, वर्कफ़्लो इवेंट, मानचित्र पहुँच",
            "वर्कफ़्लो इंजन — कॉन्फ़िगर करने योग्य JE → AE → EE → वित्त श्रृंखलाएँ",
            "सत्र प्रमाणीकरण — JWT, प्रत्येक API पर अनुमति गार्ड",
        ],
        "footer_note": "डेमो: Audit Trail — पिछले 24 घंटे फ़िल्टर करें",
        "notes": "",
    },
    {
        "type": "bullets",
        "title": "प्रस्तावित पायलट — देहरादून उत्तर मंडल",
        "hindi": "90-दिवसीय पायलट प्रस्ताव",
        "bullets": [
            "दायरा: 1 मंडल (DIV-DNN), 1–2 सक्रिय योजनाएँ, पूर्ण निर्माण + O&M + बिलिंग चक्र",
            "समयरेखा: 90 दिन — PILOT-PLAN-90-DAYS.md देखें",
            "पायलट मंडल के लिए 100% योजना डेटा डिजिटाइज़",
            "RA बिल चक्र समय 30% कम",
            "EE/SE द्वारा साप्ताहिक कार्यकारी डैशबोर्ड समीक्षा",
            "50+ क्षेत्र कर्मचारी प्रशिक्षित",
        ],
        "notes": "",
    },
    {
        "type": "table",
        "title": "अपेक्षित लाभ और ROI",
        "headers": ["क्षेत्र", "वर्तमान समस्या", "EWUMS लाभ"],
        "rows": [
            ["प्रगति रिपोर्टिंग", "मैन्युअल समेकन", "रियल-टाइम डैशबोर्ड"],
            ["भुगतान में देरी", "कागज़ MB → RA अंतराल", "डिजिटल वर्कफ़्लो"],
            ["संपत्ति रिकॉर्ड", "असंबद्ध GIS", "एकल रजिस्ट्री"],
            ["उपभोक्ता सेवा", "केवल फ़ोन शिकायतें", "पोर्टल + SLA ट्रैकिंग"],
            ["ऑडिट", "फ़ाइल पुनर्निर्माण", "अपरिवर्तनीय ऑडिट लॉग"],
        ],
        "footer_note": (
            "गुणात्मक ROI: तेज़ निर्णय, डुप्लिकेट डेटा प्रविष्टि में कमी, "
            "SFC/GoUK रिपोर्टिंग के लिए बेहतर निधि उपयोग दृश्यता।"
        ),
        "notes": "",
    },
    {
        "type": "bullets",
        "title": "PRP Geospatial Solutions — आपका कार्यान्वयन साझेदार",
        "bullets": [
            "EWUMS/S2T2R का उत्पाद स्वामी और कार्यान्वयनकर्ता",
            "उत्तराखंड-केंद्रित GIS और जल उपयोगिता डोमेन विशेषज्ञता",
            "Migration 084 मंडल संरचना पहले से कॉन्फ़िगर",
            "होस्टिंग: ewumsujs.com (प्रोडक्शन-तैयार डेमो वातावरण)",
            "सहायता मॉडल: प्रशिक्षण, डेटा माइग्रेशन, चरणबद्ध रोलआउट",
        ],
        "footer_note": "संपर्क: PRP Geospatial Solutions — [बैठक के लिए संपर्क विवरण जोड़ें]",
        "notes": "",
    },
    {
        "type": "phases",
        "title": "पायलट से आगे — राज्यव्यापी रोलआउट",
        "phases": [
            ("चरण 1 (90 दिन)", "देहरादून उत्तर पायलट"),
            ("चरण 2 (6 महीने)", "गढ़वाल सर्कल — सभी देहरादून उप-मंडल + हरिद्वार"),
            ("चरण 3 (12 महीने)", "कुमाऊँ सर्कल — नैनीताल, हल्द्वानी, रामनगर, आदि"),
            ("चरण 4", "SCADA एकीकरण, ERP लाइव पोस्टिंग, उन्नत एनालिटिक्स"),
        ],
        "notes": "",
    },
    {
        "type": "numbered",
        "title": "आज अनुरोधित निर्णय",
        "items": [
            "देहरादून उत्तर मंडल (DIV-DNN) के लिए 90-दिवसीय पायलट अनुमोदित करें",
            "पायलट टीम नामित करें — EE, 2 JE, 1 AE, 1 लेखा, 1 GIS ऑपरेटर",
            "डेटा माइग्रेशन के लिए 1–2 योजना नाम प्रदान करें (या डेमो योजनाएँ उपयोग करें)",
            "IT कनेक्टिविटी — पायलट उपयोगकर्ताओं के लिए ewumsujs.com ब्राउज़र पहुँच",
            "स्टीयरिंग समिति — PRP Geospatial के साथ पाक्षिक समीक्षा",
        ],
        "notes": "",
    },
    {
        "type": "bullets",
        "title": "प्रश्न एवं चर्चा",
        "hindi": "प्रश्न एवं चर्चा",
        "bullets": [
            "मौजूदा UJS ERP/लेखा के साथ एकीकरण",
            "जल जीवन मिशन / SFC रिपोर्टिंग संरेखण",
            "भारत के भीतर डेटा स्वामित्व और होस्टिंग",
            "हिंदी इंटरफ़ेस उपलब्धता (आंशिक — विस्तार जारी)",
        ],
        "notes": "",
    },
    {
        "type": "url",
        "title": "लाइव मंच",
        "url": "https://ewumsujs.com",
        "subtitle": "EWUMS — 17 मॉड्यूल लाइव · मंडल-स्कोप डेमो डेटा तैयार",
        "notes": "Q&A के दौरान कनेक्टिविटी अनुमति देने पर ब्राउज़र में खोलें।",
    },
    {
        "type": "closing",
        "title": "धन्यवाद — डिजिटल UJS के लिए EWUMS",
        "hindi": "धन्यवाद — डिजिटल UJS की ओर",
        "cta": (
            "पायलट अनुमोदित करें · सप्ताह 1 प्रशिक्षण निर्धारित करें · "
            "PRP औपचारिक पायलट योजना प्रस्तुत करे"
        ),
        "notes": "",
    },
]


def find_logo() -> Path | None:
    for path in LOGO_CANDIDATES:
        if path.is_file():
            return path
    return None


def set_run_font(run, size=18, bold=False, color=DARK_TEXT, name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_accent_bar(slide, prs):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def add_logo(slide, logo_path: Path, left, top, height):
    slide.shapes.add_picture(str(logo_path), left, top, height=height)


def add_header_logo(slide, prs, logo_path: Path):
    add_logo(slide, logo_path, prs.slide_width - Inches(1.55), Inches(0.18), Inches(0.42))


def add_title_logo(slide, prs, logo_path: Path):
    add_logo(slide, logo_path, prs.slide_width - Inches(3.2), Inches(0.35), Inches(1.1))
    add_logo(slide, logo_path, Inches(0.7), Inches(4.9), Inches(0.85))


def add_footer(slide, prs, logo_path: Path | None):
    if logo_path:
        add_logo(slide, logo_path, Inches(0.45), prs.slide_height - Inches(0.52), Inches(0.28))
    box = slide.shapes.add_textbox(
        Inches(0.5), prs.slide_height - Inches(0.45), prs.slide_width - Inches(1), Inches(0.35)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = FOOTER_TEXT
    set_run_font(run, size=10, color=GRAY)


def add_title_block(slide, prs, title, hindi=None, top=Inches(0.55)):
    title_box = slide.shapes.add_textbox(
        Inches(0.6), top, prs.slide_width - Inches(2.2), Inches(1.2)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run_font(run, size=28, bold=True, color=NAVY)

    if hindi:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = hindi
        set_run_font(run2, size=18, color=ORANGE)


def add_bullets(slide, prs, bullets, top=Inches(1.85), size=16):
    box = slide.shapes.add_textbox(
        Inches(0.75), top, prs.slide_width - Inches(1.5), prs.slide_height - top - Inches(0.7)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {text}"
        set_run_font(run, size=size)


def add_numbered(slide, prs, items, top=Inches(1.85), size=15):
    box = slide.shapes.add_textbox(
        Inches(0.75), top, prs.slide_width - Inches(1.5), prs.slide_height - top - Inches(0.7)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items, 1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"{i}.  {text}"
        set_run_font(run, size=size)


def add_table_slide(slide, prs, headers, rows, top=Inches(1.75)):
    cols = len(headers)
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        cols,
        Inches(0.6),
        top,
        prs.slide_width - Inches(1.2),
        Inches(0.4 * (len(rows) + 1)),
    )
    table = table_shape.table
    col_width = (prs.slide_width - Inches(1.2)) / cols
    for c in range(cols):
        table.columns[c].width = int(col_width)

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                set_run_font(run, size=12, bold=True, color=WHITE)

    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run_font(run, size=11, color=DARK_TEXT)


def add_footer_note(slide, prs, text):
    box = slide.shapes.add_textbox(
        Inches(0.6), prs.slide_height - Inches(0.95), prs.slide_width - Inches(1.2), Inches(0.45)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run_font(run, size=11, color=GRAY, bold=False)


def add_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def build_slide(prs, data, logo_path: Path | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_accent_bar(slide, prs)

    stype = data["type"]

    if stype == "title":
        add_footer(slide, prs, logo_path)
        if logo_path:
            add_title_logo(slide, prs, logo_path)

        band = slide.shapes.add_shape(1, Inches(0), Inches(0.12), prs.slide_width, Inches(3.2))
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY
        band.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.6), prs.slide_width - Inches(4.0), Inches(1.6)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = data["title"]
        set_run_font(run, size=32, bold=True, color=WHITE)

        sub_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.5), prs.slide_width - Inches(1.4), Inches(0.6)
        )
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = data["subtitle"]
        set_run_font(run2, size=20, color=ORANGE)

        tag_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(3.6), prs.slide_width - Inches(1.4), Inches(0.5)
        )
        tf3 = tag_box.text_frame
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        run3.text = data.get("tagline", "")
        set_run_font(run3, size=14, color=GRAY)

        url_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(4.2), prs.slide_width - Inches(1.4), Inches(0.4)
        )
        tf4 = url_box.text_frame
        p4 = tf4.paragraphs[0]
        run4 = p4.add_run()
        run4.text = "https://ewumsujs.com"
        set_run_font(run4, size=16, bold=True, color=ORANGE)

    elif stype == "closing":
        add_footer(slide, prs, logo_path)
        if logo_path:
            add_header_logo(slide, prs, logo_path)

        band = slide.shapes.add_shape(1, Inches(0), Inches(0.12), prs.slide_width, Inches(4.5))
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY
        band.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.2), prs.slide_width - Inches(1.4), Inches(1.2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = data["title"]
        set_run_font(run, size=30, bold=True, color=WHITE)

        if data.get("hindi"):
            p2 = tf.add_paragraph()
            run2 = p2.add_run()
            run2.text = data["hindi"]
            set_run_font(run2, size=22, color=ORANGE)

        cta_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(3.0), prs.slide_width - Inches(1.4), Inches(0.8)
        )
        tf3 = cta_box.text_frame
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        run3.text = data.get("cta", "")
        set_run_font(run3, size=14, color=WHITE)

        url_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(4.8), prs.slide_width - Inches(1.4), Inches(0.4)
        )
        tf4 = url_box.text_frame
        p4 = tf4.paragraphs[0]
        run4 = p4.add_run()
        run4.text = "https://ewumsujs.com"
        set_run_font(run4, size=18, bold=True, color=ORANGE)

    else:
        add_footer(slide, prs, logo_path)
        if logo_path:
            add_header_logo(slide, prs, logo_path)

        if stype == "url":
            add_title_block(slide, prs, data["title"])
            url_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(2.5), prs.slide_width - Inches(1.2), Inches(1.0)
            )
            tf = url_box.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = data["url"]
            set_run_font(run, size=36, bold=True, color=ORANGE)
            if data.get("subtitle"):
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                run2 = p2.add_run()
                run2.text = data["subtitle"]
                set_run_font(run2, size=16, color=GRAY)

        elif stype == "demo":
            add_title_block(slide, prs, data["title"], data.get("hindi"))
            lines = [
                f"लॉगिन: {data['login']}",
                f"नेविगेट: {data['navigate']}",
                "",
                "दिखाएँ:",
            ] + [f"  • {s}" for s in data["show"]]
            if data.get("script"):
                lines += ["", f'स्क्रिप्ट: "{data["script"]}"']
            add_bullets(slide, prs, lines, top=Inches(1.75), size=14)

        elif stype == "bullets":
            add_title_block(slide, prs, data["title"], data.get("hindi"))
            add_bullets(slide, prs, data["bullets"])
            if data.get("footer_note"):
                add_footer_note(slide, prs, data["footer_note"])

        elif stype == "numbered":
            add_title_block(slide, prs, data["title"], data.get("hindi"))
            add_numbered(slide, prs, data["items"])
            if data.get("footer_note"):
                add_footer_note(slide, prs, data["footer_note"])

        elif stype == "table":
            add_title_block(slide, prs, data["title"], data.get("hindi"))
            add_table_slide(slide, prs, data["headers"], data["rows"])
            if data.get("footer_note"):
                add_footer_note(slide, prs, data["footer_note"])

        elif stype == "sections":
            add_title_block(slide, prs, data["title"], data.get("hindi"))
            box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.85), prs.slide_width - Inches(1.5), Inches(3.5)
            )
            tf = box.text_frame
            tf.word_wrap = True
            for i, (label, desc) in enumerate(data["sections"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(10)
                run_label = p.add_run()
                run_label.text = f"{label}: "
                set_run_font(run_label, size=16, bold=True, color=NAVY)
                run_desc = p.add_run()
                run_desc.text = desc
                set_run_font(run_desc, size=16, color=DARK_TEXT)
            if data.get("footer_note"):
                add_footer_note(slide, prs, data["footer_note"])

        elif stype == "phases":
            add_title_block(slide, prs, data["title"])
            box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.85), prs.slide_width - Inches(1.5), Inches(3.5)
            )
            tf = box.text_frame
            for i, (phase, desc) in enumerate(data["phases"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(12)
                run_p = p.add_run()
                run_p.text = f"{phase}: "
                set_run_font(run_p, size=17, bold=True, color=ORANGE)
                run_d = p.add_run()
                run_d.text = desc
                set_run_font(run_d, size=17, color=DARK_TEXT)

    add_notes(slide, data.get("notes", ""))


def main():
    logo_path = find_logo()
    if logo_path:
        print(f"Using logo: {logo_path}")
    else:
        print("Warning: prp-logo.png not found — generating without logo")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in SLIDES:
        build_slide(prs, slide_data, logo_path)

    prs.save(OUTPUT)
    print(f"Created {OUTPUT} with {len(SLIDES)} slides (Hindi only)")


if __name__ == "__main__":
    main()
