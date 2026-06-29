#!/usr/bin/env python3
"""Generate EWUMS UJS department presentation (.pptx) from PRESENTATION-OUTLINE.md content."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# PRP brand colors
NAVY = RGBColor(0x0F, 0x17, 0x2A)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
FOOTER_TEXT = "PRP Geospatial Solutions | ewumsujs.com"

OUTPUT = Path(__file__).parent / "EWUMS-UJS-Presentation.pptx"

SLIDES = [
    {
        "type": "title",
        "title": "Integrated Water Utility Management\nfor Uttarakhand Jal Sansthan",
        "subtitle": "एकीकृत जल उपयोगिता प्रबंधन प्रणाली — UJS",
        "tagline": "EWUMS — Enterprises Water Utility Management System (S2T2R)",
        "notes": (
            "Good morning. Today we present EWUMS — a single digital platform covering "
            "planning, construction, GIS, operations, billing, and executive oversight for UJS."
        ),
    },
    {
        "type": "bullets",
        "title": "Why UJS Needs a Unified Platform",
        "hindi": "चुनौती — बिखरी हुई जानकारी और manual प्रक्रियाएँ",
        "bullets": [
            "Schemes tracked in spreadsheets, files, and separate GIS tools",
            "Delayed visibility on physical vs financial progress",
            "Division-wise data silos (Garhwal / Kumaon circles)",
            "Consumer billing, complaints, and O&M not linked to assets on map",
            "Audit and approval trails difficult to reconstruct",
        ],
        "notes": (
            "Officers spend time reconciling data instead of managing schemes. "
            "EWUMS addresses this end-to-end."
        ),
    },
    {
        "type": "bullets",
        "title": "EWUMS — One Platform, Full Lifecycle",
        "hindi": "समाधान — योजना से O&M तक एक ही प्रणाली",
        "bullets": [
            "21 integrated modules — 17 live today on ewumsujs.com",
            "DPR → Construction → MB → RA Bill → Final Bill → GIS handover → O&M → Billing",
            "Division-scoped access aligned to UJS field structure (Migration 084 divisions)",
            "Web GIS with district boundaries (Dehradun, Haridwar, statewide)",
            "Mobile billing and consumer portal (Jal Mitra)",
        ],
        "notes": "",
    },
    {
        "type": "table",
        "title": "Built for UJS Organisation Structure",
        "headers": ["Level", "Roles", "Platform Access"],
        "rows": [
            ["State HQ (Dehradun)", "MD, CGM, CE", "Statewide dashboard & all divisions"],
            ["Circle", "SE Garhwal / Kumaon", "Circle-wide schemes"],
            ["Division", "EE, AE, JE", "Division-scoped projects & GIS"],
            ["Field", "Contractor, Accounts", "Construction workflow & payments"],
            ["Citizen", "Consumer", "Jal Mitra portal — bills, complaints"],
        ],
        "footer_note": "Reference divisions: Dehradun North (DIV-DNN), Haridwar Rural (DIV-HRR), "
        "Karanprayag (DIV-KPG), Nainital (DIV-NTL)",
        "notes": "",
    },
    {
        "type": "numbered",
        "title": "Stage 1–2: DPR & Approval",
        "hindi": "योजना पाइपलाइन (12 चरण)",
        "items": [
            "DPR proposal initiation (Division EE/JE)",
            "HQ review → TAC → Secretariat → Sanction",
            "Tender published → ready for project",
        ],
        "footer_note": "Screens: /dpr-planning · /workflows",
        "notes": (
            "Walk through the 12-stage DPR approval pipeline — from division proposal "
            "to tender publication."
        ),
    },
    {
        "type": "bullets",
        "title": "Stage 3–4: Land & Project",
        "hindi": "भूमि अधिग्रहण → परियोजना पंजीकरण",
        "bullets": [
            "GIS pipeline trace + parcel clearances (/land-acquisition)",
            "HQ / Super Admin registers project from tender-published DPR (/projects)",
            "Division staff logins auto-provisioned (JE, AE, EE, Accounts)",
            "Milestones and budget auto-linked to approved DPR",
        ],
        "notes": (
            "Only Super Admin creates schemes. Division is tagged at creation — "
            "all Karanprayag staff see the same project end-to-end."
        ),
    },
    {
        "type": "bullets",
        "title": "Stage 5–6: GIS & Construction",
        "hindi": "मानचित्र पर निर्माण — हर रुपये का हिसाब",
        "bullets": [
            "Map assets: pipe, reservoir, FHTC (/map)",
            "Daily DPR → MB → BOQ reconciliation → RA Bill",
            "JE → AE → EE approval chain with full audit trail",
            "Physical % and financial % on construction dashboard",
        ],
        "footer_note": "Screen: /projects/{id}/construction",
        "notes": "",
    },
    {
        "type": "bullets",
        "title": "Stage 7: Handover to O&M",
        "hindi": "कमीशनिंग व O&M सौंपना",
        "bullets": [
            "Completion checklist — MB, GIS, FHTC, commissioning certificates",
            "Handover certificate + asset register",
            "Scheme becomes operational for billing and consumer service",
        ],
        "footer_note": "Screen: /om#handover",
        "notes": "",
    },
    {
        "type": "numbered",
        "title": "Stage 8–9: Consumers & Revenue",
        "hindi": "FHTC पंजीकरण → बिलिंग → राजस्व",
        "items": [
            "Register consumer (FHTC + mobile) linked to scheme",
            "Meter reading via mobile billing (/mobile-billing)",
            "Bill generation → collection and arrears tracking",
            "15.12 reports · revenue GIS analytics by geography",
        ],
        "footer_note": "Screens: /billing · /mobile-billing",
        "notes": "",
    },
    {
        "type": "table",
        "title": "Stage 10: Citizen Complaints (Jal Mitra Only)",
        "hindi": "शिकायत — नागरिक दर्ज करे, विभाग निपटाए",
        "headers": ["Who", "Action"],
        "rows": [
            ["Consumer", "Registers complaint via Jal Mitra (/portal)"],
            ["EE", "Sees ticket → assigns JE"],
            ["JE", "Resolves in field"],
            ["System", "Email alert on assignment"],
        ],
        "footer_note": (
            "IMPORTANT: Staff do NOT register complaints — only consumers do. "
            "Flow: Ticket Generated → Assigned → Resolution → Feedback → Closed"
        ),
        "notes": "",
    },
    {
        "type": "bullets",
        "title": "Stage 11–12: Command Center & Audit",
        "hindi": "कमांड सेंटर + ऑडिट ट्रेल",
        "bullets": [
            "KPIs: assets, progress, revenue, complaints, SLA",
            "Division or statewide view (/dashboard)",
            "Every user action logged — workflows, map access, approvals",
            "Audit Trail page — filter by user, division, last 24 hours",
        ],
        "footer_note": "Screens: /dashboard · /admin/audit",
        "notes": "",
    },
    {
        "type": "table",
        "title": "Live Demo — Karanprayag Flow (15 min)",
        "hindi": "लाइव डेमो — करणप्रयाग प्रवाह",
        "headers": ["Step", "Who", "Action"],
        "rows": [
            ["1", "Admin", "Create scheme (Karanprayag division)"],
            ["2", "Admin", "Register consumer (FHTC + mobile)"],
            ["3", "Consumer", "Jal Mitra — lodge complaint"],
            ["4", "EE", "Complaints inbox — assign JE"],
            ["5", "All", "Notification log + email alert"],
        ],
        "footer_note": "Login: admin@egip.local / Admin@123 · ee.kpg@egip.local / EE@123",
        "notes": "Rehearse full proposal-to-revenue + complaint flow before department meeting.",
    },
    {
        "type": "bullets",
        "title": "Proposed Pilot — Dehradun North Division",
        "hindi": "90-दिवसीय पायलट प्रस्ताव",
        "bullets": [
            "Scope: 1 division (DIV-DNN), 1–2 active schemes, full construction + O&M + billing cycle",
            "Timeline: 90 days — see PILOT-PLAN-90-DAYS.md",
            "100% scheme data digitised for pilot division",
            "RA bill cycle time reduced by 30%",
            "Executive dashboard reviewed weekly by EE/SE",
            "50+ field staff trained",
        ],
        "notes": "",
    },
    {
        "type": "table",
        "title": "Expected Benefits & ROI",
        "headers": ["Area", "Current Pain", "EWUMS Benefit"],
        "rows": [
            ["Progress reporting", "Manual consolidation", "Real-time dashboard"],
            ["Payment delays", "Paper MB → RA lag", "Digital workflow"],
            ["Asset records", "Disconnected GIS", "Single registry"],
            ["Consumer service", "Phone-only complaints", "Portal + SLA tracking"],
            ["Audit", "File reconstruction", "Immutable audit log"],
        ],
        "footer_note": (
            "Qualitative ROI: Faster decision-making, reduced duplicate data entry, "
            "improved fund utilisation visibility for SFC/GoUK reporting."
        ),
        "notes": "",
    },
    {
        "type": "bullets",
        "title": "PRP Geospatial Solutions — Your Implementation Partner",
        "bullets": [
            "Product owner and implementer of EWUMS/S2T2R",
            "Uttarakhand-focused GIS and water utility domain expertise",
            "Migration 084 division structure already configured",
            "Hosting: ewumsujs.com (production-ready demo environment)",
            "Support model: training, data migration, phased rollout",
        ],
        "footer_note": "Contact: PRP Geospatial Solutions — [Add contact details for meeting]",
        "notes": "",
    },
    {
        "type": "phases",
        "title": "Beyond Pilot — Statewide Rollout",
        "phases": [
            ("Phase 1 (90 days)", "Dehradun North pilot"),
            ("Phase 2 (6 months)", "Garhwal Circle — all Dehradun sub-divisions + Haridwar"),
            ("Phase 3 (12 months)", "Kumaon Circle — Nainital, Haldwani, Ramnagar, etc."),
            ("Phase 4", "SCADA integration, ERP live posting, advanced analytics"),
        ],
        "notes": "",
    },
    {
        "type": "numbered",
        "title": "Decisions Requested Today",
        "items": [
            "Approve 90-day pilot for Dehradun North Division (DIV-DNN)",
            "Nominate pilot team — EE, 2 JE, 1 AE, 1 Accounts, 1 GIS operator",
            "Provide 1–2 scheme names for data migration (or use demo schemes)",
            "IT connectivity — browser access to ewumsujs.com for pilot users",
            "Steering committee — fortnightly review with PRP Geospatial",
        ],
        "notes": "",
    },
    {
        "type": "bullets",
        "title": "Questions & Discussion",
        "hindi": "प्रश्न एवं चर्चा",
        "bullets": [
            "Integration with existing UJS ERP/accounting",
            "Jal Jeevan Mission / SFC reporting alignment",
            "Data ownership and hosting within India",
            "Hindi interface availability (partial — expanding)",
        ],
        "notes": "",
    },
    {
        "type": "url",
        "title": "Live Platform",
        "url": "https://ewumsujs.com",
        "subtitle": "EWUMS — 17 modules live · Division-scoped demo data ready",
        "notes": "Open in browser during Q&A if connectivity allows.",
    },
    {
        "type": "closing",
        "title": "Thank You — EWUMS for a Digital UJS",
        "hindi": "धन्यवाद — डिजिटल UJS की ओर",
        "cta": "Approve pilot · Schedule training Week 1 · PRP to submit formal pilot plan",
        "notes": "",
    },
]


def set_run_font(run, size=18, bold=False, color=DARK_TEXT, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_accent_bar(slide, prs):
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0),
        Inches(0),
        prs.slide_width,
        Inches(0.12),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def add_footer(slide, prs):
    box = slide.shapes.add_textbox(
        Inches(0.5), prs.slide_height - Inches(0.45), prs.slide_width - Inches(1), Inches(0.35)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = FOOTER_TEXT
    set_run_font(run, size=10, color=GRAY)


def add_title_block(slide, prs, title, hindi=None, top=Inches(0.55)):
    title_box = slide.shapes.add_textbox(
        Inches(0.6), top, prs.slide_width - Inches(1.2), Inches(1.2)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run_font(run, size=28, bold=True, color=NAVY)

    if hindi:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = hindi
        set_run_font(run2, size=18, color=ORANGE)


def add_bullets(slide, prs, bullets, top=Inches(1.85), size=16):
    box = slide.shapes.add_textbox(
        Inches(0.75), top, prs.slide_width - Inches(1.5), prs.slide_height - top - Inches(0.7)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {text}"
        set_run_font(run, size=size)


def add_numbered(slide, prs, items, top=Inches(1.85), size=15):
    box = slide.shapes.add_textbox(
        Inches(0.75), top, prs.slide_width - Inches(1.5), prs.slide_height - top - Inches(0.7)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items, 1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"{i}.  {text}"
        set_run_font(run, size=size)


def add_table_slide(slide, prs, headers, rows, top=Inches(1.75)):
    cols = len(headers)
    table_shape = slide.shapes.add_table(
        len(rows) + 1, cols, Inches(0.6), top, prs.slide_width - Inches(1.2), Inches(0.4 * (len(rows) + 1))
    )
    table = table_shape.table
    col_width = (prs.slide_width - Inches(1.2)) / cols
    for c in range(cols):
        table.columns[c].width = int(col_width)

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                set_run_font(run, size=12, bold=True, color=WHITE)

    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run_font(run, size=11, color=DARK_TEXT)


def add_footer_note(slide, prs, text):
    box = slide.shapes.add_textbox(
        Inches(0.6), prs.slide_height - Inches(0.95), prs.slide_width - Inches(1.2), Inches(0.45)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run_font(run, size=11, color=GRAY, bold=False)


def add_notes(slide, text):
    if text:
        notes = slide.notes_slide
        notes.notes_text_frame.text = text


def build_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_accent_bar(slide, prs)
    add_footer(slide, prs)

    stype = data["type"]

    if stype == "title":
        # Navy header band
        band = slide.shapes.add_shape(
            1, Inches(0), Inches(0.12), prs.slide_width, Inches(3.2)
        )
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY
        band.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.6), prs.slide_width - Inches(1.4), Inches(1.6)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = data["title"]
        set_run_font(run, size=32, bold=True, color=WHITE)

        sub_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.5), prs.slide_width - Inches(1.4), Inches(0.6)
        )
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = data["subtitle"]
        set_run_font(run2, size=20, color=ORANGE)

        tag_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(3.6), prs.slide_width - Inches(1.4), Inches(0.5)
        )
        tf3 = tag_box.text_frame
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        run3.text = data.get("tagline", "")
        set_run_font(run3, size=14, color=GRAY)

        url_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(4.2), prs.slide_width - Inches(1.4), Inches(0.4)
        )
        tf4 = url_box.text_frame
        p4 = tf4.paragraphs[0]
        run4 = p4.add_run()
        run4.text = "https://ewumsujs.com"
        set_run_font(run4, size=16, bold=True, color=ORANGE)

    elif stype == "closing":
        band = slide.shapes.add_shape(
            1, Inches(0), Inches(0.12), prs.slide_width, Inches(4.5)
        )
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY
        band.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.2), prs.slide_width - Inches(1.4), Inches(1.2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = data["title"]
        set_run_font(run, size=30, bold=True, color=WHITE)

        if data.get("hindi"):
            p2 = tf.add_paragraph()
            run2 = p2.add_run()
            run2.text = data["hindi"]
            set_run_font(run2, size=22, color=ORANGE)

        cta_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(3.0), prs.slide_width - Inches(1.4), Inches(0.8)
        )
        tf3 = cta_box.text_frame
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        run3.text = data.get("cta", "")
        set_run_font(run3, size=14, color=WHITE)

        url_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(4.8), prs.slide_width - Inches(1.4), Inches(0.4)
        )
        tf4 = url_box.text_frame
        p4 = tf4.paragraphs[0]
        run4 = p4.add_run()
        run4.text = "https://ewumsujs.com"
        set_run_font(run4, size=18, bold=True, color=ORANGE)

    elif stype == "url":
        add_title_block(slide, prs, data["title"])
        url_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(2.5), prs.slide_width - Inches(1.2), Inches(1.0)
        )
        tf = url_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = data["url"]
        set_run_font(run, size=36, bold=True, color=ORANGE)
        if data.get("subtitle"):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = data["subtitle"]
            set_run_font(run2, size=16, color=GRAY)

    elif stype == "demo":
        add_title_block(slide, prs, data["title"], data.get("hindi"))
        lines = [
            f"Login: {data['login']}",
            f"Navigate: {data['navigate']}",
            "",
            "Show:",
        ] + [f"  • {s}" for s in data["show"]]
        if data.get("script"):
            lines += ["", f'Script: "{data["script"]}"']
        add_bullets(slide, prs, lines, top=Inches(1.75), size=14)

    elif stype == "bullets":
        add_title_block(slide, prs, data["title"], data.get("hindi"))
        add_bullets(slide, prs, data["bullets"])
        if data.get("footer_note"):
            add_footer_note(slide, prs, data["footer_note"])

    elif stype == "numbered":
        add_title_block(slide, prs, data["title"], data.get("hindi"))
        add_numbered(slide, prs, data["items"])
        if data.get("footer_note"):
            add_footer_note(slide, prs, data["footer_note"])

    elif stype == "table":
        add_title_block(slide, prs, data["title"], data.get("hindi"))
        add_table_slide(slide, prs, data["headers"], data["rows"])
        if data.get("footer_note"):
            add_footer_note(slide, prs, data["footer_note"])

    elif stype == "sections":
        add_title_block(slide, prs, data["title"], data.get("hindi"))
        box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.85), prs.slide_width - Inches(1.5), Inches(3.5)
        )
        tf = box.text_frame
        tf.word_wrap = True
        for i, (label, desc) in enumerate(data["sections"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            run_label = p.add_run()
            run_label.text = f"{label}: "
            set_run_font(run_label, size=16, bold=True, color=NAVY)
            run_desc = p.add_run()
            run_desc.text = desc
            set_run_font(run_desc, size=16, color=DARK_TEXT)
        if data.get("footer_note"):
            add_footer_note(slide, prs, data["footer_note"])

    elif stype == "phases":
        add_title_block(slide, prs, data["title"])
        box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.85), prs.slide_width - Inches(1.5), Inches(3.5)
        )
        tf = box.text_frame
        for i, (phase, desc) in enumerate(data["phases"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(12)
            run_p = p.add_run()
            run_p.text = f"{phase}: "
            set_run_font(run_p, size=17, bold=True, color=ORANGE)
            run_d = p.add_run()
            run_d.text = desc
            set_run_font(run_d, size=17, color=DARK_TEXT)

    add_notes(slide, data.get("notes", ""))


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in SLIDES:
        build_slide(prs, slide_data)

    prs.save(OUTPUT)
    print(f"Created {OUTPUT} with {len(SLIDES)} slides")


if __name__ == "__main__":
    main()
